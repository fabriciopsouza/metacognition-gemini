#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""doc_intake.py — parse determinístico + chunk + manifesto (ADR-029).

Ingestão de documentos OFFLINE e AGNÓSTICA para alimentar o discovery: lê
pdf/docx/xlsx/pptx/md/txt, extrai texto de forma determinística, fatia em chunks
com overlap e emite um MANIFESTO JSON auditável. SEM embeddings, SEM rede
(RAG-vetorial é ADR futuro com dependência externa declarada — decisão travada #4).

Princípios:
  - Determinístico: mesmo input -> mesmo manifesto (sem timestamps por default;
    chunking por limites de parágrafo; sha256 por arquivo e por chunk).
  - Degrada com segurança: formato cujo parser opcional não está instalado vira
    status "skipped" no manifesto — NUNCA derruba o processo (os demais seguem).
  - Stdlib-first: .md/.txt funcionam sempre; pdf/docx/xlsx/pptx usam libs opcionais
    (pypdf, python-docx, openpyxl, python-pptx) importadas sob demanda.

Uso:
  python doc_intake.py <arquivo-ou-pasta> [--out manifest.json]
                       [--chunk-chars 1200] [--overlap 150] [--with-text]

Saída: manifesto JSON em stdout (ou --out). `--with-text` inclui o texto dos chunks
(default: só hash + offsets, para manifesto leve e estável).
"""

import argparse
import hashlib
import json
import os
import sys

SUPPORTED = {".pdf", ".docx", ".xlsx", ".pptx", ".md", ".txt", ".markdown", ".text"}


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


# --- Extractors por formato. Cada um retorna (texto, extractor, status). -------
# status: "ok" | "skipped: <razão>" | "error: <razão>". Texto vazio em skip/error.

def _extract_textlike(path: str):
    with open(path, "rb") as fh:
        raw = fh.read()
    # Decodifica tolerante (utf-8 -> latin-1 fallback) para não quebrar em BOM/encoding.
    try:
        text = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = raw.decode("latin-1", errors="replace")
    return text, "stdlib:text", "ok"


def _extract_pdf(path: str):
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except ImportError:
            return "", "pdf", "skipped: requer 'pypdf' (pip install pypdf)"
    try:
        reader = PdfReader(path)
        parts = [(page.extract_text() or "") for page in reader.pages]
        return "\n\n".join(parts), "pypdf", "ok"
    except Exception as exc:  # parser de PDF falha em arquivos corrompidos
        return "", "pypdf", f"error: {type(exc).__name__}: {exc}"


def _extract_docx(path: str):
    try:
        import docx  # type: ignore  (python-docx)
    except ImportError:
        return "", "docx", "skipped: requer 'python-docx' (pip install python-docx)"
    try:
        document = docx.Document(path)
        parts = [p.text for p in document.paragraphs]
        return "\n".join(parts), "python-docx", "ok"
    except Exception as exc:
        return "", "python-docx", f"error: {type(exc).__name__}: {exc}"


def _extract_xlsx(path: str):
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return "", "xlsx", "skipped: requer 'openpyxl' (pip install openpyxl)"
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"# sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    lines.append("\t".join(cells))
        return "\n".join(lines), "openpyxl", "ok"
    except Exception as exc:
        return "", "openpyxl", f"error: {type(exc).__name__}: {exc}"


def _extract_pptx(path: str):
    try:
        from pptx import Presentation  # type: ignore  (python-pptx)
    except ImportError:
        return "", "pptx", "skipped: requer 'python-pptx' (pip install python-pptx)"
    try:
        prs = Presentation(path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"# slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(run.text for run in para.runs)
                        if txt:
                            lines.append(txt)
        return "\n".join(lines), "python-pptx", "ok"
    except Exception as exc:
        return "", "python-pptx", f"error: {type(exc).__name__}: {exc}"


EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".md": _extract_textlike,
    ".markdown": _extract_textlike,
    ".txt": _extract_textlike,
    ".text": _extract_textlike,
}


def extract(path: str):
    ext = os.path.splitext(path)[1].lower()
    fn = EXTRACTORS.get(ext)
    if fn is None:
        return "", ext or "?", f"skipped: formato não suportado ({ext or 'sem extensão'})"
    return fn(path)


# --- Chunking determinístico ---------------------------------------------------

def chunk_text(text: str, chunk_chars: int, overlap: int):
    """Fatia o texto em janelas de ~chunk_chars com overlap, preferindo cortar em
    fronteira de parágrafo/linha próxima ao limite. Determinístico e sem perda."""
    if chunk_chars <= 0:
        raise ValueError("chunk_chars deve ser > 0")
    if overlap < 0 or overlap >= chunk_chars:
        raise ValueError("overlap deve estar em [0, chunk_chars)")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    n = len(text)
    if n == 0:
        return []
    chunks = []
    start = 0
    while start < n:
        end = min(start + chunk_chars, n)
        # Tenta recuar até a última quebra de parágrafo/linha dentro da janela
        # (não recua mais que 40% do chunk, para não gerar chunks minúsculos).
        if end < n:
            window_floor = start + int(chunk_chars * 0.6)
            cut = text.rfind("\n\n", window_floor, end)
            if cut == -1:
                cut = text.rfind("\n", window_floor, end)
            if cut == -1:
                cut = text.rfind(" ", window_floor, end)
            if cut > start:
                end = cut
        piece = text[start:end]
        chunks.append((start, end, piece))
        if end >= n:
            break
        start = max(end - overlap, start + 1)
    return chunks


# --- Manifesto -----------------------------------------------------------------

def build_file_entry(path: str, chunk_chars: int, overlap: int, with_text: bool, id_base: str = None):
    # id_base: prefixo do chunk-id. Default = path relativo (não basename) para garantir
    # UNICIDADE entre arquivos de mesmo nome em subpastas distintas (MÉDIO-1 qa-critic).
    if id_base is None:
        id_base = path.replace("\\", "/")
    try:
        with open(path, "rb") as fh:
            raw = fh.read()
        size = len(raw)
        file_sha = sha256_bytes(raw)
    except OSError as exc:
        # Ramo de erro mantém o MESMO schema dos entries normais (sentinelas), para o
        # consumidor poder iterar sem KeyError (BAIXO-1 qa-critic).
        return {
            "path": path.replace("\\", "/"),
            "format": os.path.splitext(path)[1].lower(),
            "bytes": 0,
            "sha256": "",
            "extractor": "",
            "status": f"error: {type(exc).__name__}: {exc}",
            "n_chars": 0,
            "n_chunks": 0,
            "chunks": [],
        }

    text, extractor, status = extract(path)
    entry = {
        "path": path.replace("\\", "/"),
        "format": os.path.splitext(path)[1].lower(),
        "bytes": size,
        "sha256": file_sha,
        "extractor": extractor,
        "status": status,
        "n_chars": len(text),
        "n_chunks": 0,
        "chunks": [],
    }
    if status != "ok" or not text.strip():
        if status == "ok" and not text.strip():
            entry["status"] = "ok: vazio (sem texto extraível)"
        return entry

    pieces = chunk_text(text, chunk_chars, overlap)
    for ordinal, (cstart, cend, ptext) in enumerate(pieces):
        chunk = {
            "id": f"{id_base}#{ordinal:04d}",
            "ordinal": ordinal,
            "char_start": cstart,
            "char_end": cend,
            "n_chars": len(ptext),
            "sha256": sha256_text(ptext),
        }
        if with_text:
            chunk["text"] = ptext
        entry["chunks"].append(chunk)
    entry["n_chunks"] = len(pieces)
    return entry


def collect_paths(target: str):
    if os.path.isfile(target):
        return [target]
    if os.path.isdir(target):
        found = []
        for root, _dirs, files in os.walk(target):
            for name in files:
                if os.path.splitext(name)[1].lower() in SUPPORTED:
                    found.append(os.path.join(root, name))
        return sorted(found)  # ordem determinística
    return []


def build_manifest(target: str, chunk_chars: int, overlap: int, with_text: bool):
    paths = collect_paths(target)
    # id_base = path relativo ao target (pasta) -> chunk-ids únicos mesmo com basenames
    # repetidos em subpastas. Para target arquivo único, usa o basename.
    base_dir = target if os.path.isdir(target) else os.path.dirname(target)

    def _id_base(p):
        try:
            rel = os.path.relpath(p, base_dir)
        except ValueError:  # drives diferentes no Windows
            rel = os.path.basename(p)
        return rel.replace("\\", "/")

    files = [build_file_entry(p, chunk_chars, overlap, with_text, _id_base(p)) for p in paths]
    return {
        "manifest_version": 1,
        "tool": "doc_intake.py",
        "params": {"chunk_chars": chunk_chars, "overlap": overlap, "with_text": with_text},
        "embeddings": False,  # explícito: ingestão é só parse+chunk (decisão #4)
        "summary": {
            "n_files": len(files),
            "n_ok": sum(1 for f in files if f["status"].startswith("ok")),
            "n_skipped": sum(1 for f in files if f["status"].startswith("skipped")),
            "n_error": sum(1 for f in files if f["status"].startswith("error")),
            "n_chunks": sum(f["n_chunks"] for f in files),
        },
        "files": files,
    }


def main(argv=None):
    ap = argparse.ArgumentParser(description="Parse determinístico + chunk + manifesto (ADR-029).")
    ap.add_argument("target", help="arquivo ou pasta a ingerir")
    ap.add_argument("--out", default="", help="caminho do manifesto JSON (default: stdout)")
    ap.add_argument("--chunk-chars", type=int, default=1200, help="tamanho-alvo do chunk em chars")
    ap.add_argument("--overlap", type=int, default=150, help="overlap entre chunks em chars")
    ap.add_argument("--with-text", action="store_true", help="inclui o texto dos chunks no manifesto")
    args = ap.parse_args(argv)

    if not os.path.exists(args.target):
        sys.stderr.write(f"[doc_intake] alvo inexistente: {args.target}\n")
        return 2

    manifest = build_manifest(args.target, args.chunk_chars, args.overlap, args.with_text)
    payload = json.dumps(manifest, ensure_ascii=False, indent=2, sort_keys=False)
    if args.out:
        with open(args.out, "w", encoding="utf-8") as fh:
            fh.write(payload)
        sys.stderr.write(
            f"[doc_intake] manifesto: {args.out} — {manifest['summary']['n_files']} arquivo(s), "
            f"{manifest['summary']['n_chunks']} chunk(s), {manifest['summary']['n_skipped']} pulado(s).\n"
        )
    else:
        sys.stdout.write(payload + "\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
