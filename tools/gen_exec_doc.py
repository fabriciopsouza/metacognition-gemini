#!/usr/bin/env python3
"""gen_exec_doc.py — elaboração de DOCUMENTOS premium, flexível por TIPO (ADR-050). PREMIUM-only.

Gera o documento que cada situação pede — proposta/orçamento (custo + trade-offs + alternativas +
aprovação), **POP/SOP**, **manual de operação**, **guia de configuração**, **plano de manutenção**,
etc. — em **md/docx/pptx/pdf**. **QUAL documento cada situação exige é definido pelo discovery /
explorer / briefing + PMO** (papel deles; este gerador NÃO decide — apenas RENDERIZA o que a spec
declara). Agnóstico: a ESTRUTURA vem do template/spec; o CONTEÚDO vem do projeto.

NÃO hardcoda seções: renderiza **as seções `##` que a spec tiver**. Se a spec declarar
`<!-- required: A; B; C -->`, valida que essas seções existem (a "definição" do tipo é da spec/template,
não deste tool). ANTI-FABRICAÇÃO: campo vazio/placeholder vira **"NÃO PREENCHIDO"** — nunca inventado
(custo, passo de POP, parâmetro de config — qualquer campo).

Formatos: **md** sempre (stdlib). **docx** (python-docx), **pptx** (python-pptx), **pdf** (reportlab) se
a lib existir; ausente → pula o formato (degrada, não falha).

Uso:
    python tools/gen_exec_doc.py <spec.md> --out-dir <dir> [--formats md,docx,pptx,pdf]
Exit 0 se gerou (md no mínimo) e as seções `required` (se declaradas) existem; 1 caso contrário.
"""
import argparse
import datetime
import os
import re
import sys
import unicodedata

PLACEHOLDER = re.compile(r"^\s*(<.*>|preencher|a preencher|tbd|todo|\.\.\.|\?|n/d)\s*$", re.IGNORECASE)
REQUIRED_RX = re.compile(r"<!--\s*required:\s*(.+?)\s*-->", re.IGNORECASE)
NAO_PREENCHIDO = "NÃO PREENCHIDO"

# Convenção de ENTREGA (ADR-050, emenda): output/<datestamp>-<label>/ com subpastas por tipo.
# Torna a entrega navegável pelo usuário leigo. Cada formato cai na sua subpasta.
FORMAT_SUBDIR = {"md": "docs", "docx": "docs", "pdf": "docs", "pptx": "apresentacao"}
DELIVERY_SUBDIRS = {
    "codigo": "scripts e código-fonte da entrega",
    "docs": "documentos (md/docx/pdf): relatórios, manuais, POPs, decisão",
    "apresentacao": "apresentações executivas (pptx)",
    "dados": "planilhas e dados de entrada/saída",
}


def slugify(s):
    """ASCII-safe, minúsculo, hifenizado — nome de pasta seguro cross-platform (ç→c, ã→a)."""
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode("ascii")
    s = re.sub(r"\(.*?\)", "", s)                       # remove parênteses (ex.: "(premium, ADR-050)")
    s = re.sub(r"[^\w\s-]", "", s).strip().lower()
    return re.sub(r"[\s_-]+", "-", s) or "entrega"


def parse_title(text):
    m = re.search(r"^#\s+(.*)$", text, re.MULTILINE)
    return m.group(1).strip() if m else "Documento"


def parse_sections(text):
    """Retorna lista [(titulo, conteudo)] das seções ## (preserva ordem)."""
    out, cur, buf = [], None, []
    for ln in text.splitlines():
        h = re.match(r"^#{2,6}\s+(.*)$", ln)
        if h:
            if cur is not None:
                out.append((cur, "\n".join(buf).strip()))
            cur, buf = h.group(1).strip(), []
        elif cur is not None:
            buf.append(ln)
    if cur is not None:
        out.append((cur, "\n".join(buf).strip()))
    return out


def _norm(s):
    return re.sub(r"\s+", " ", s.strip().lower())


def required_sections(text):
    """Seções obrigatórias declaradas pela spec/template (`<!-- required: A; B -->`). [] se não declara."""
    m = REQUIRED_RX.search(text)
    if not m:
        return []
    return [r.strip() for r in re.split(r"[;,]", m.group(1)) if r.strip()]


def value_or_nao_preenchido(content):
    """Anti-fabricação: vazio/placeholder → NÃO PREENCHIDO. Nunca inventa."""
    if not content or PLACEHOLDER.match(content) or not re.search(r"[a-z0-9]", content, re.IGNORECASE):
        return NAO_PREENCHIDO
    lines = [l for l in content.splitlines() if l.strip()]
    if lines and all(PLACEHOLDER.match(re.sub(r"^[-*]\s*", "", l)) for l in lines):
        return NAO_PREENCHIDO
    return content


def validate(text):
    """(ok, problemas). Falta de seção `required` declarada = problema. Título ausente = problema."""
    problems = []
    if "#" not in text or not parse_title(text):
        problems.append("sem título (# ...)")
    present = {_norm(t) for t, _ in parse_sections(text)}
    for req in required_sections(text):
        if _norm(req) not in present:
            problems.append(f"seção obrigatória (required) ausente: {req}")
    ok = not any("ausente" in p or "sem título" in p for p in problems)
    return ok, problems


def render_markdown(text):
    title = parse_title(text)
    L = [f"# {title}", ""]
    for t, content in parse_sections(text):
        if REQUIRED_RX.search(t):  # não renderiza o próprio marcador como seção
            continue
        L.append(f"## {t}")
        L.append(value_or_nao_preenchido(content))
        L.append("")
    L.append("> Gerado por `tools/gen_exec_doc.py` (premium, ADR-050). Tipo e seções definidos pela spec "
             "(discovery/explorer/PMO decidem qual documento a situação pede). Campos vazios = **NÃO "
             "PREENCHIDO** — nunca fabricados.")
    return "\n".join(L) + "\n"


def _chunk(s, n):
    """Quebra texto em pedaços de até n chars sem chunk vazio à frente (paginação, não truncagem)."""
    out, buf = [], ""
    for line in s.splitlines() or [""]:
        while len(line) > n:                       # linha sozinha maior que o limite
            if buf:
                out.append(buf); buf = ""
            out.append(line[:n]); line = line[n:]
        if not buf:
            buf = line
        elif len(buf) + 1 + len(line) <= n:
            buf += "\n" + line
        else:
            out.append(buf); buf = line
    if buf or not out:
        out.append(buf)
    return out


def _wrap(line, n):
    out = []
    while len(line) > n:
        out.append(line[:n]); line = line[n:]
    out.append(line)
    return out


def _dst(out_dir, subdirs, fmt, name):
    """Caminho de saída roteado por subpasta (modo --deliver) ou plano (default), sem overwrite cego."""
    d = os.path.join(out_dir, subdirs[fmt]) if subdirs and fmt in subdirs else out_dir
    os.makedirs(d, exist_ok=True)
    p = os.path.join(d, name)
    if os.path.exists(p):                           # nunca sobrescreve em silêncio: bump -2/-3...
        stem, ext = os.path.splitext(name)
        i = 2
        while os.path.exists(os.path.join(d, f"{stem}-{i}{ext}")):
            i += 1
        p = os.path.join(d, f"{stem}-{i}{ext}")
    return p


def export(text, out_dir, formats, subdirs=None, basename="documento"):
    os.makedirs(out_dir, exist_ok=True)
    title = parse_title(text)
    secs = [(t, value_or_nao_preenchido(c)) for t, c in parse_sections(text) if not REQUIRED_RX.search(t)]
    written, skipped = [], []
    if "md" in formats:
        p = _dst(out_dir, subdirs, "md", f"{basename}.md")
        open(p, "w", encoding="utf-8").write(render_markdown(text))
        written.append(p)
    if "docx" in formats:
        try:
            import docx  # type: ignore
            d = docx.Document(); d.add_heading(title, 0)
            for t, c in secs:
                d.add_heading(t, level=1); d.add_paragraph(c)
            p = _dst(out_dir, subdirs, "docx", f"{basename}.docx"); d.save(p); written.append(p)
        except Exception as e:
            skipped.append(f"docx ({e}; instale python-docx)")
    if "pptx" in formats:
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation()
            t0 = prs.slides.add_slide(prs.slide_layouts[0]); t0.shapes.title.text = title
            for t, c in secs:
                parts = _chunk(c, 700)             # pagina em vez de truncar silenciosamente
                for i, part in enumerate(parts):
                    s = prs.slides.add_slide(prs.slide_layouts[1])
                    s.shapes.title.text = t if len(parts) == 1 else f"{t} ({i + 1}/{len(parts)})"
                    s.placeholders[1].text = part
            p = _dst(out_dir, subdirs, "pptx", f"{basename}.pptx"); prs.save(p); written.append(p)
        except Exception as e:
            skipped.append(f"pptx ({e}; instale python-pptx)")
    if "pdf" in formats:
        try:
            from reportlab.lib.pagesizes import A4  # type: ignore
            from reportlab.pdfgen import canvas  # type: ignore
            p = _dst(out_dir, subdirs, "pdf", f"{basename}.pdf"); c = canvas.Canvas(p, pagesize=A4); y = 800

            def _line(s, x, size, bold):
                nonlocal y
                if y < 60:
                    c.showPage(); y = 800
                c.setFont("Helvetica-Bold" if bold else "Helvetica", size)
                c.drawString(x, y, s)
                y -= size + 4

            _line(title[:95], 40, 15, True); y -= 8
            for t, content in secs:
                _line(t[:95], 40, 12, True)
                for line in content.splitlines():   # sem corte em 8 linhas: pagina tudo
                    for seg in _wrap(line, 95):
                        _line(seg, 48, 10, False)
                y -= 8
            c.save(); written.append(p)
        except Exception as e:
            skipped.append(f"pdf ({e}; instale reportlab ou converta o .docx)")
    return written, skipped


def main(argv):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("--out-dir", default="saida-doc")
    ap.add_argument("--formats", default="md,docx,pptx,pdf")
    ap.add_argument("--deliver", metavar="LABEL",
                    help="monta entrega em <root>/<datestamp>-<LABEL>/ com subpastas + índice navegável")
    ap.add_argument("--root", default="output", help="raiz das entregas (default: output/)")
    args = ap.parse_args(argv[1:])
    if not os.path.isfile(args.spec):
        print(f"spec inexistente: {args.spec}", file=sys.stderr)
        return 2
    text = open(args.spec, encoding="utf-8-sig").read()
    ok, problems = validate(text)
    for p in problems:
        print(("FALTA: " if "ausente" in p or "sem título" in p else "AVISO: ") + p)
    if not ok:
        print("RESULTADO: FAIL (documento incompleto — seção obrigatória declarada ausente)")
        return 1
    formats = [f.strip() for f in args.formats.split(",") if f.strip()]
    title = parse_title(text)
    if args.deliver:
        stamp = datetime.datetime.now().strftime("%Y-%m-%dT%H%M%S")
        out_dir = os.path.join(args.root, f"{stamp}-{slugify(args.deliver)}")
        for sub in DELIVERY_SUBDIRS:                      # cria todas as subpastas-padrão da entrega
            os.makedirs(os.path.join(out_dir, sub), exist_ok=True)
        written, skipped = export(text, out_dir, formats, subdirs=FORMAT_SUBDIR, basename=slugify(title))
    else:
        out_dir, written, skipped = args.out_dir, *export(text, args.out_dir, formats)
    for w in written:
        print(f"  gerado: {w}")
    for s in skipped:
        print(f"  pulado: {s}")
    if args.deliver:
        try:
            sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
            import make_index                              # baseline; índice navegável + handoff
            n, idx = make_index.generate(out_dir, title=title)
            for p in idx:
                print(f"  índice: {p}")
            print(f"ENTREGA: {out_dir} ({n} arquivo(s); abra index.html)")
        except Exception as e:
            print(f"  AVISO: índice não gerado ({e}); rode tools/make_index.py {out_dir}")
    print("RESULTADO: PASS (documento gerado; tipo definido pela spec; campos vazios = NÃO PREENCHIDO)")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
